"""Build docs/Safety_and_Trust_Engine_US.pptx from the EU deck.

Copies the existing presentation and rewrites regulatory framing for the
US pack (NIST AI RMF + federal MRM). Architecture / toolkit / demo slides
stay; honesty caveats (guidance / supervisory) replace "legal obligation".
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "docs" / "Safety_and_Trust_Engine.pptx"
DST = ROOT / "docs" / "Safety_and_Trust_Engine_US.pptx"

# slide index (1-based) -> list of (old, new) exact shape-text matches
SLIDE_REPLACEMENTS: dict[int, list[tuple[str, str]]] = {
    1: [
        (
            "garak · AgentDojo · PyRIT    |    EU AI Act · DORA · FCA Operational Resilience",
            "garak · AgentDojo · PyRIT    |    NIST AI RMF · Federal MRM (SR 26-02)",
        ),
    ],
    2: [
        (
            "Adversarial testing became a legal obligation",
            "Adversarial testing is now an exam-ready expectation",
        ),
        ("17 Jan 2025", "26 Jul 2024"),
        ("DORA", "NIST AI 600-1"),
        (
            "Fully applicable across EU financial entities \uf0e0 https://www.eiopa.europa.eu/digital-operational-resilience-act-dora_en",
            "Generative AI Profile — red-team as pre-deploy MEASURE \uf0e0 https://nvlpubs.nist.gov/nistpubs/ai/NIST.AI.600-1.pdf",
        ),
        ("31 Mar 2025", "17 Apr 2026"),
        ("FCA PS21/3", "SR 26-02"),
        (
            "Operational-resilience transition period closed \uf0e0 https://www.fca.org.uk/publications/policy-statements/ps21-3-building-operational-resilience",
            "Revised interagency MRM guidance (Fed / OCC / FDIC) \uf0e0 https://www.federalreserve.gov/supervisionreg/srletters/SR2602.pdf",
        ),
        ("Art. 15 / 55", "MEASURE 2.6–2.7"),
        ("EU AI Act", "NIST AI RMF"),
        (
            "Robustness, cybersecurity & documented red-teaming \uf0e0\nhttps://artificialintelligenceact.eu/article/15/",
            "Misuse, security & resilience evaluated and documented \uf0e0\nhttps://www.nist.gov/itl/ai-risk-management-framework",
        ),
        ("Three regulators, one shared demand", "Two US lenses, one shared demand"),
        (
            "Adversarial testing that is repeatable, evidenced, and remediated. Most teams already run the attacks — a notebook, a one-off run, a screenshot. None of that is evidence. The gap this engine closes is turning red-teaming into something you can prove, on every commit.",
            "Adversarial testing that is repeatable, evidenced, and remediated. NIST is voluntary guidance; federal MRM is supervisory (GenAI often by analogy under SR 26-02). The gap this engine closes is turning red-teaming into something you can prove, on every commit — with binding honesty in the artifact.",
        ),
    ],
    4: [
        ("EU AI Act — Article 15 (and 55)", "NIST AI RMF — MEASURE (and AI 600-1)"),
        ("Art. 15(1)", "MEASURE 2.6"),
        (
            "Appropriate accuracy, robustness and cybersecurity, consistent across the lifecycle.",
            "AI system evaluated for potential for misuse and abuse.",
        ),
        ("Art. 15(5)", "MEASURE 2.7"),
        (
            "Resilience against unauthorised third parties exploiting vulnerabilities — incl. data poisoning and adversarial inputs.",
            "AI system security and resilience are evaluated and documented.",
        ),
        ("Art. 55(1)(a)", "AI 600-1"),
        (
            "GPAI with systemic risk must conduct and document adversarial testing (red-teaming).",
            "Generative AI Profile: pre-deployment red-teaming for GenAI information-security risks.",
        ),
        (
            "Engineering translation:  “document” is the hook — our evidence artifact is that documentation.",
            "Engineering translation:  binding = guidance — voluntary NIST, not certification. The artifact is the documentation.",
        ),
    ],
    5: [
        ("DORA — testing and third-party pillars", "Federal MRM — validation and monitoring"),
        ("Art. 24–25", "SR 26-02 validation"),
        (
            "Risk-based resilience testing: vulnerability assessments, independent testers, all critical tools tested ≥ annually.",
            "Independent validation / effective challenge: outcomes analysis via adversarial testing.",
        ),
        ("Art. 26–27", "SR 26-02 monitoring"),
        (
            "Threat-led penetration testing (TLPT) ≥ every 3 years, simulating real-world threat actors.",
            "Ongoing monitoring: repeatable adversarial gate in CI with numeric ASR tolerances.",
        ),
        ("Art. 28", "Third-party model"),
        (
            "ICT third-party risk — AWS Bedrock is an ICT third party, in scope of testing and the register.",
            "Vendor / hosted model provider recorded in the artifact and kept in testing scope.",
        ),
        (
            "Engineering translation:  a nightly CI run is continuous assurance — not a substitute for formal TLPT.",
            "Engineering translation:  continuous CI assurance ≠ formal exam exercise; GenAI is supervisory analogy under SR 26-02.",
        ),
    ],
    6: [
        (
            "FCA PS21/3 — operational resilience",
            "Federal MRM — documentation & decision package",
        ),
        ("Important business service", "Model-risk file"),
        (
            "An agent is a dependency of an important business service whose disruption can cause intolerable harm.",
            "The gate’s JSON + Markdown artifact is the validation-and-documentation evidence a model-risk reviewer asks for.",
        ),
        ("Impact tolerance", "Numeric ASR gate"),
        (
            "The maximum tolerable disruption — tested under severe-but-plausible scenarios.",
            "Maximum acceptable attack-success rate per category — same mechanic, US vocabulary.",
        ),
        ("Self-assessment", "Documentation"),
        (
            "A written record evidencing resilience and remediation.",
            "Satisfied by the consolidated evidence artifact (binding: supervisory).",
        ),
        (
            "Engineering translation:  impact tolerance is just a maximum acceptable attack-success rate.",
            "Engineering translation:  keep binding honest — supervisory expectation, not a statute with penalties.",
        ),
    ],
    7: [
        (
            "In code today: EU AI Act · DORA · FCA.  NIST AI RMF and the MIT AI Risk Repository are the natural next lenses.",
            "In code today: --regimes us encodes NIST AI RMF · Federal MRM. Default pack remains eu_uk; combine with --regimes eu_uk,us.",
        ),
    ],
    11: [
        ("Impact tolerance, made numeric", "Exam-ready tolerance, made numeric"),
        (
            "FCA impact tolerance\n=\nmax attack-success rate",
            "SR 26-02 monitoring\n=\nmax attack-success rate",
        ),
    ],
    12: [
        (
            'Control(\n  regulation="EU AI Act",\n  ref="Art. 15(5)",\n  stages=(GARAK, AGENTDOJO, PYRIT),\n  categories=("prompt_injection",\n    "tool_injection", "jailbreak"),\n)',
            'Control(\n  regulation="NIST AI RMF",\n  ref="MEASURE 2.7",\n  binding="guidance",\n  stages=(GARAK, AGENTDOJO, PYRIT),\n  categories=("prompt_injection",\n    "tool_injection", "jailbreak"),\n)',
        ),
    ],
    17: [
        ("Red-teaming is now regulated", "Red-teaming is now exam-ready"),
        (
            "EU AI Act, DORA and FCA all demand repeatable, documented adversarial testing — one control, three vocabularies.",
            "NIST MEASURE and federal MRM both reward repeatable, documented adversarial testing — guidance and supervisory lenses, not silent certification.",
        ),
        (
            "Resources   README · docs/REGULATORY_RESEARCH.md · docs/pipeline.svg · docs/safety_trust_engine_cicd_pipeline.svg · .github/workflows/safety-trust.yml · docs/HANDOVER.md",
            "Resources   README · docs/US_COMPLIANCE.md · docs/REGULATORY_RESEARCH.md §4.5 · docs/pipeline.svg · .github/workflows/safety-trust.yml",
        ),
    ],
}


def _set_shape_text(shape, new: str) -> None:
    tf = shape.text_frame
    first = True
    for para in tf.paragraphs:
        if first:
            if para.runs:
                para.runs[0].text = new
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.text = new
            first = False
        else:
            for run in para.runs:
                run.text = ""
            if not para.runs:
                para.text = ""


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"missing source deck: {SRC}")
    shutil.copy2(SRC, DST)
    prs = Presentation(str(DST))
    hits = 0
    missing: list[str] = []
    for idx, pairs in SLIDE_REPLACEMENTS.items():
        slide = prs.slides[idx - 1]
        shapes = [
            sh for sh in slide.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()
        ]
        for old, new in pairs:
            matched = False
            for shape in shapes:
                if shape.text_frame.text == old:
                    _set_shape_text(shape, new)
                    hits += 1
                    matched = True
                    break
            if not matched:
                missing.append(f"slide {idx}: {old[:70]!r}")
    prs.save(str(DST))
    print(f"wrote {DST} ({hits} replacements)")
    if missing:
        print("UNMATCHED:")
        for m in missing:
            print(" ", m)
        # Debug: print near-misses for unmatched
        for idx, pairs in SLIDE_REPLACEMENTS.items():
            slide = prs.slides[idx - 1]
            for old, _new in pairs:
                if any(
                    sh.has_text_frame and sh.text_frame.text == old
                    for sh in slide.shapes
                ):
                    continue
                # still missing after failed run — show candidates containing a keyword
                key = old.split("\n", 1)[0][:40]
                for sh in slide.shapes:
                    if sh.has_text_frame and key[:20] in sh.text_frame.text:
                        print(f"  candidate slide {idx}: {sh.text_frame.text!r}")
                        break



if __name__ == "__main__":
    main()
